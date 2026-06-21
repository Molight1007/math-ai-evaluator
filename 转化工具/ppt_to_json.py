"""
PowerPoint (.pptx / .ppt) → JSON 题目转换工具。

将 PowerPoint 格式的数学题集转换为评测器可用的 JSON 格式。
逐页提取文本，识别题号+选项模式，合并为完整题目。

用法:
    python 转化工具/ppt_to_json.py <ppt路径> [-o 输出.json] [--max N]

输出 JSON 格式:
[
  {
    "id": "ppt_001",
    "question": "题目完整内容（含选项）",
    "domain": "",
    "reference_answer": ""
  }
]
"""
import argparse
import json
import os
import re
import sys

try:
    from pptx import Presentation
except ImportError:
    print("请先安装 python-pptx: pip install python-pptx")
    sys.exit(1)


def _extract_slide_text(slide) -> str:
    """
    从单页幻灯片中提取所有文本内容。

    参数:
        slide: python-pptx Slide 对象

    返回:
        该页所有文本（换行分隔）
    """
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return "\n".join(texts)


def _is_problem_start(text: str) -> bool:
    """判断文本是否以题号开头"""
    return bool(re.match(r'^\d+[\.\、\s]', text.strip()))


def _is_answer_slide(text: str) -> bool:
    """判断当前页是否为答案页（含答案/解析关键词）"""
    return bool(re.search(
        r'(?:参考答案|答案解析|答案[：:]|习题答案|标准答案)',
        text,
    ))


def convert_ppt(ppt_path: str, max_problems: int = 0) -> list[dict]:
    """
    将 PowerPoint 文件转换为题目列表。

    逐页提取文本后按题号分割题目，自动跳过答案页。
    支持 .pptx 格式，.ppt 格式需先另存为 .pptx。

    参数:
        ppt_path: PowerPoint 文件路径
        max_problems: 最大提取题目数（0 表示全部）

    返回:
        题目字典列表，每项含 id/question/domain/reference_answer
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"文件不存在: {ppt_path}")

    ext = os.path.splitext(ppt_path)[1].lower()
    if ext == ".ppt":
        # 旧版 .ppt 尝试用 python-pptx 打开
        try:
            prs = Presentation(ppt_path)
        except Exception:
            raise ValueError(
                "无法直接读取旧版 .ppt 文件。"
                "请用 PowerPoint 打开后另存为 .pptx 格式，再重新导入。"
            )
    else:
        prs = Presentation(ppt_path)

    problems = []
    in_answer_section = False
    problem_counter = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_text = _extract_slide_text(slide)
        if not slide_text:
            continue

        # 检测答案页
        if _is_answer_slide(slide_text):
            in_answer_section = True
            continue
        if in_answer_section:
            continue

        # 按行分割，逐行识别题目
        lines = slide_text.split("\n")
        for line in lines:
            line = line.strip()
            if not line:
                continue

            if _is_problem_start(line):
                problem_counter += 1
                match = re.match(r'^(\d+)', line)
                num = int(match.group(1)) if match else problem_counter
                question_text = re.sub(r'^\d+[\.\、\s]+', '', line).strip()

                problems.append({
                    "id": f"ppt_{problem_counter:04d}",
                    "question": question_text,
                    "domain": "",
                    "reference_answer": "",
                })
            elif problems:
                # 续行：拼接到上一题（可能是选项或题干续文）
                problems[-1]["question"] += " " + line

    # 后处理：清理空白
    for p in problems:
        p["question"] = re.sub(r'\s+', ' ', p["question"]).strip()

    print(f"从 PPT 解析出 {len(problems)} 道题目 (共 {len(prs.slides)} 页)")

    if max_problems > 0:
        problems = problems[:max_problems]
        print(f"限制为前 {max_problems} 道")

    return problems


def main() -> None:
    """命令行入口"""
    parser = argparse.ArgumentParser(
        description="PowerPoint (.pptx/.ppt) → JSON 题目转换器"
    )
    parser.add_argument("ppt", help="PowerPoint 文件路径 (.pptx 或 .ppt)")
    parser.add_argument("-o", "--output", default=None, help="输出 JSON 文件路径")
    parser.add_argument("--max", type=int, default=0, help="最多提取题目数 (0=全部)")
    args = parser.parse_args()

    ppt_path = args.ppt
    if not os.path.exists(ppt_path):
        print(f"文件不存在: {ppt_path}")
        sys.exit(1)

    problems = convert_ppt(ppt_path, max_problems=args.max)

    # 默认输出路径
    if args.output is None:
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        output_dir = os.path.join(base_dir, "测试结果", "原本问题")
        os.makedirs(output_dir, exist_ok=True)
        ppt_name = os.path.splitext(os.path.basename(ppt_path))[0]
        args.output = os.path.join(output_dir, f"{ppt_name}.json")

    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(problems, f, ensure_ascii=False, indent=2)

    print(f"\n已保存到: {args.output}")
    print(f"共 {len(problems)} 道题目")


if __name__ == "__main__":
    main()
