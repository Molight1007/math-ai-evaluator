"""
答案文档提取器 — 从多种格式的答案文件中提取「题目-答案」对。

支持格式:
- .pptx / .ppt (PowerPoint): 按幻灯片逐页提取
- .docx (Word): 按段落提取，识别题号+答案模式
- .txt / .md: 纯文本/Markdown，按行解析
- .pdf: 用 pdfplumber 提取文本后解析
- .csv: 用 pandas 读 CSV，自动识别题目列和答案列
- .xlsx: 用 openpyxl 读取，按行列解析
- .json: 从 JSON 中提取 question/answer 字段对

输出统一格式:
[
  {
    "index": 1,              # 序号
    "question_text": "...",  # 题干（可能为空，如果答案文档只有答案）
    "answer_text": "...",    # 答案内容
    "source_page": 1         # 来源页码/幻灯片号
  }
]
"""
import csv
import json
import logging
import os
import re

logger = logging.getLogger(__name__)

# 支持的文件扩展名集合（用于快速成员检查）
SUPPORTED_ANSWER_EXTS: set = {
    ".pptx", ".ppt", ".docx", ".txt", ".md",
    ".pdf", ".csv", ".xlsx", ".json",
}


def extract_from_pptx(pptx_path: str) -> list[dict]:
    """
    从 PowerPoint 文件中提取答案对。
    策略：逐页提取所有文本，用正则匹配「题目+答案」模式。
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")

    prs = Presentation(pptx_path)
    pairs = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        # 提取当前页所有文本
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)

        if not texts:
            continue

        full_text = "\n".join(texts)

        # 尝试从当前页中提取题目-答案对
        page_pairs = _parse_answer_pairs(full_text, slide_idx)
        pairs.extend(page_pairs)

    logger.info(f"从 PPT 提取了 {len(pairs)} 个答案对 (共 {len(prs.slides)} 页)")
    return pairs


def extract_from_docx(docx_path: str) -> list[dict]:
    """
    从 Word 文档中提取答案对。
    策略：按段落读取，识别题号模式 + 答案标记。
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")

    doc = Document(docx_path)
    pairs = []
    current_page = 1  # Word 没有明确的页码概念，用段落序号近似

    # 先收集所有段落
    all_paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            all_paragraphs.append(text)

    # 也读取表格中的内容
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                all_paragraphs.append(" | ".join(row_texts))

    full_text = "\n".join(all_paragraphs)
    pairs = _parse_answer_pairs(full_text, current_page)

    logger.info(f"从 Word 提取了 {len(pairs)} 个答案对 (共 {len(all_paragraphs)} 个段落)")
    return pairs


def extract_from_txt(txt_path: str) -> list[dict]:
    """从纯文本文件提取答案对"""
    with open(txt_path, "r", encoding="utf-8") as f:
        content = f.read()

    pairs = _parse_answer_pairs(content, 1)
    logger.info(f"从文本文件提取了 {len(pairs)} 个答案对")
    return pairs


def extract_answers(filepath: str) -> list[dict]:
    """
    统一入口：根据文件扩展名自动选择提取器。

    参数:
        filepath: 答案文件路径

    返回:
        list[dict]，每个 dict 包含 index, question_text, answer_text, source_page

    异常:
        ValueError: 不支持的文件格式
    """
    ext = os.path.splitext(filepath)[1].lower()

    # 按扩展名分发到对应的提取函数
    _dispatchers = {
        ".pptx": extract_from_pptx,
        ".ppt": extract_from_ppt,
        ".docx": extract_from_docx,
        ".txt": extract_from_txt,
        ".md": extract_from_md,
        ".pdf": extract_from_pdf,
        ".csv": extract_from_csv,
        ".xlsx": extract_from_xlsx,
        ".json": extract_from_json,
    }

    extract_fn = _dispatchers.get(ext)
    if extract_fn is None:
        supported = ", ".join(sorted(SUPPORTED_ANSWER_EXTS))
        raise ValueError(f"不支持的答案文件格式: {ext}（支持 {supported}）")

    pairs = extract_fn(filepath)

    # 重新编号，确保 index 连续
    for i, pair in enumerate(pairs):
        pair["index"] = i + 1

    return pairs


def extract_from_ppt(ppt_path: str) -> list[dict]:
    """
    从旧版 PowerPoint (.ppt) 文件中提取答案对。

    旧版 .ppt 是 OLE 二进制格式，python-pptx 不支持直接读取。
    优先尝试用 python-pptx 打开（部分 .ppt 内部实为 .pptx 格式），
    失败则提示用户将文件另存为 .pptx。

    参数:
        ppt_path: .ppt 文件路径

    返回:
        提取的答案对列表
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")

    try:
        # 尝试用 python-pptx 打开（部分 .ppt 内部兼容 .pptx 格式）
        prs = Presentation(ppt_path)
        pairs = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if not texts:
                continue
            full_text = "\n".join(texts)
            page_pairs = _parse_answer_pairs(full_text, slide_idx)
            pairs.extend(page_pairs)
        logger.info(f"从 .ppt 提取了 {len(pairs)} 个答案对 (共 {len(prs.slides)} 页)")
        return pairs
    except Exception:
        # .ppt 无法直接用 python-pptx 打开，提示用户转换格式
        raise ValueError(
            "无法直接读取旧版 .ppt 文件。"
            "请用 PowerPoint 打开后另存为 .pptx 格式，再重新导入。"
        )


def extract_from_md(md_path: str) -> list[dict]:
    """
    从 Markdown 文件中提取答案对。

    按标题（# 开头）分段后逐段解析题目-答案对，
    复用 _parse_answer_pairs 的通用解析逻辑。

    参数:
        md_path: .md 文件路径

    返回:
        提取的答案对列表
    """
    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()

    pairs = _parse_answer_pairs(content, 1)
    logger.info(f"从 Markdown 文件提取了 {len(pairs)} 个答案对")
    return pairs


def extract_from_pdf(pdf_path: str) -> list[dict]:
    """
    从 PDF 文件中提取答案对。

    使用 pdfplumber 逐页提取文本，然后按页解析题目-答案对。
    每页独立解析，避免跨页边界匹配问题。

    参数:
        pdf_path: .pdf 文件路径

    返回:
        提取的答案对列表
    """
    try:
        import pdfplumber
    except ImportError:
        raise ImportError("请安装 pdfplumber: pip install pdfplumber")

    pairs = []
    with pdfplumber.open(pdf_path) as pdf:
        for page_idx, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if not text:
                continue
            text = text.strip()
            if not text:
                continue
            # 逐页独立解析，source_page 标记实际页码
            page_pairs = _parse_answer_pairs(text, page_idx)
            pairs.extend(page_pairs)

    logger.info(f"从 PDF 提取了 {len(pairs)} 个答案对 (共 {len(pdf.pages)} 页)")
    return pairs


def extract_from_csv(csv_path: str) -> list[dict]:
    """
    从 CSV 文件中提取答案对。

    自动识别题目列和答案列的策略：
    1. 先查找列名中包含「题目/问题/question」的列作为题目列
    2. 再查找列名中包含「答案/answer」的列作为答案列
    3. 若未识别到，默认第一列为题目列、最后一列为答案列

    参数:
        csv_path: .csv 文件路径

    返回:
        提取的答案对列表
    """
    import pandas as pd

    df = pd.read_csv(csv_path, encoding="utf-8")
    if df.empty:
        logger.warning("CSV 文件为空")
        return []

    columns = [str(c).strip().lower() for c in df.columns]

    # 自动识别题目列：列名包含 question/题目/问题 等关键词
    question_col_idx = None
    answer_col_idx = None
    for i, col in enumerate(columns):
        if question_col_idx is None and _is_column_like(col, ("question", "题目", "问题", "题干", "problem")):
            question_col_idx = i
        if answer_col_idx is None and _is_column_like(col, ("answer", "答案", "解答", "解析", "solution")):
            answer_col_idx = i

    # 未识别到则用默认策略
    if question_col_idx is None:
        question_col_idx = 0
    if answer_col_idx is None:
        answer_col_idx = len(df.columns) - 1

    pairs = []
    for idx, row in df.iterrows():
        question_text = str(row.iloc[question_col_idx]) if pd.notna(row.iloc[question_col_idx]) else ""
        answer_text = str(row.iloc[answer_col_idx]) if pd.notna(row.iloc[answer_col_idx]) else ""
        # 跳过全空行
        if not question_text.strip() and not answer_text.strip():
            continue
        pairs.append({
            "index": idx + 1,
            "question_text": question_text.strip(),
            "answer_text": answer_text.strip(),
            "source_page": idx + 1,
        })

    logger.info(f"从 CSV 提取了 {len(pairs)} 个答案对")
    return pairs


def _is_column_like(col_name: str, keywords: tuple[str, ...]) -> bool:
    """判断列名是否包含指定关键词（不区分大小写）"""
    col_lower = col_name.lower()
    return any(kw.lower() in col_lower for kw in keywords)


def extract_from_xlsx(xlsx_path: str) -> list[dict]:
    """
    从 Excel (.xlsx) 文件中提取答案对。

    与 CSV 提取策略一致：自动识别题目列和答案列。
    使用 openpyxl 读取，支持 .xlsx 格式。

    参数:
        xlsx_path: .xlsx 文件路径

    返回:
        提取的答案对列表
    """
    import pandas as pd

    df = pd.read_excel(xlsx_path, engine="openpyxl")
    if df.empty:
        logger.warning("Excel 文件为空")
        return []

    columns = [str(c).strip().lower() for c in df.columns]

    question_col_idx = None
    answer_col_idx = None
    for i, col in enumerate(columns):
        if question_col_idx is None and _is_column_like(col, ("question", "题目", "问题", "题干", "problem")):
            question_col_idx = i
        if answer_col_idx is None and _is_column_like(col, ("answer", "答案", "解答", "解析", "solution")):
            answer_col_idx = i

    if question_col_idx is None:
        question_col_idx = 0
    if answer_col_idx is None:
        answer_col_idx = len(df.columns) - 1

    pairs = []
    for idx, row in df.iterrows():
        question_text = str(row.iloc[question_col_idx]) if pd.notna(row.iloc[question_col_idx]) else ""
        answer_text = str(row.iloc[answer_col_idx]) if pd.notna(row.iloc[answer_col_idx]) else ""
        if not question_text.strip() and not answer_text.strip():
            continue
        pairs.append({
            "index": idx + 1,
            "question_text": question_text.strip(),
            "answer_text": answer_text.strip(),
            "source_page": idx + 1,
        })

    logger.info(f"从 Excel 提取了 {len(pairs)} 个答案对")
    return pairs


def extract_from_json(json_path: str) -> list[dict]:
    """
    从 JSON 文件中提取答案对。

    支持两种 JSON 结构：
    1. 对象数组：[{"question": "...", "answer": "..."}, ...]
    2. 单个对象：{"questions": [...], "answers": [...]}

    自动识别 question/answer 字段的常见别名。

    参数:
        json_path: .json 文件路径

    返回:
        提取的答案对列表
    """
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    pairs = []

    # 结构1：对象数组
    if isinstance(data, list):
        for idx, item in enumerate(data):
            if not isinstance(item, dict):
                continue
            q = _find_value_by_keys(item, ("question", "题目", "题干", "problem", "q"))
            a = _find_value_by_keys(item, ("answer", "答案", "解答", "解析", "solution", "a"))
            if not q and not a:
                continue
            pairs.append({
                "index": idx + 1,
                "question_text": str(q).strip() if q else "",
                "answer_text": str(a).strip() if a else "",
                "source_page": idx + 1,
            })

    # 结构2：单个对象，包含 questions/answers 数组
    elif isinstance(data, dict):
        questions = data.get("questions") or data.get("题目列表") or []
        answers = data.get("answers") or data.get("答案列表") or []
        max_len = max(len(questions), len(answers))
        for i in range(max_len):
            q = questions[i] if i < len(questions) else ""
            a = answers[i] if i < len(answers) else ""
            if not q and not a:
                continue
            pairs.append({
                "index": i + 1,
                "question_text": str(q).strip() if q else "",
                "answer_text": str(a).strip() if a else "",
                "source_page": i + 1,
            })

    logger.info(f"从 JSON 提取了 {len(pairs)} 个答案对")
    return pairs


def _find_value_by_keys(data: dict, keys: tuple[str, ...]) -> str | None:
    """
    在字典中按优先级查找第一个存在的 key 对应的值。

    参数:
        data: 待查找的字典
        keys: 候选 key 元组，按优先级从高到低排列

    返回:
        找到的值（字符串），未找到返回 None
    """
    for key in keys:
        value = data.get(key)
        if value is not None and str(value).strip():
            return str(value).strip()
    return None


def _parse_answer_pairs(text: str, page: int) -> list[dict]:
    """
    从文本中解析「题目-答案」对。
    支持多种常见格式：

    格式1: 题号 + 题干 + 答案
        1. 求函数 f(x)=x² 的导数
        答案: 2x

    格式2: 题号 + 答案（无题干）
        1. B
        2. C

    格式3: 表格形式
        1 | 求导数 | 2x

    格式4: 答案列表
        一、1. A  2. B  3. C

    格式5: 题号. 题干 （多行）答案: xxx
    """
    pairs = []

    # ---- 策略1：按题号分段 ----
    # 匹配各种题号格式: "1.", "1、", "1)", "(1)", "第1题", "Q1.", "问题1"
    problem_pattern = re.compile(
        r'(?:^|\n)\s*'
        r'(?:第\s*)?(\d+)\s*[\.\、\)）题]?\s*'
        r'(?![\d\.\、\)）])',  # 后面不能紧跟数字（排除 1.2.3 这种情况）
        re.MULTILINE
    )

    # 找到所有题号的位置
    matches = list(problem_pattern.finditer(text))

    if len(matches) >= 1:
        for i, match in enumerate(matches):
            num = int(match.group(1))
            start = match.end()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
            segment = text[start:end].strip()

            # 在 segment 中分离题干和答案
            question_text, answer_text = _split_question_answer(segment)

            pairs.append({
                "index": num,
                "question_text": question_text,
                "answer_text": answer_text,
                "source_page": page,
            })

        return pairs

    # ---- 策略2：按「答案」关键词分割 ----
    answer_pattern = re.compile(
        r'(?:^|\n)\s*(?:答案|解答|解析|参考答案)[：:\s]*(.*?)(?=\n\s*(?:答案|解答|解析|参考答案|\d+[\.\、\)）])|\Z)',
        re.DOTALL | re.MULTILINE
    )
    answer_matches = list(answer_pattern.finditer(text))

    if answer_matches:
        idx = 1
        for m in answer_matches:
            answer_text = m.group(1).strip()[:500]  # 截断过长答案
            pairs.append({
                "index": idx,
                "question_text": "",
                "answer_text": answer_text,
                "source_page": page,
            })
            idx += 1
        return pairs

    # ---- 策略3：按行分割，每行一个答案 ----
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    # 如果每行都很短（<100字符），可能每行是一个独立答案
    if lines and all(len(l) < 200 for l in lines):
        for idx, line in enumerate(lines, 1):
            # 尝试分离题号和内容
            q, a = _split_question_answer(line)
            pairs.append({
                "index": idx,
                "question_text": q,
                "answer_text": a or line,
                "source_page": page,
            })

    return pairs


def _split_question_answer(segment: str) -> tuple[str, str]:
    """
    从一段文本中分离题干和答案。

    按优先级尝试多种匹配模式：
    1. 明确的「答案: xxx」标记
    2. 以「故选/因此/所以」结尾的行作为答案
    3. 单独字母 A-D（选择题答案）
    4. 短文本整体当答案

    参数:
        segment: 待分离的文本段

    返回:
        (question_text, answer_text) 元组
    """
    segment = segment.strip()

    # 模式1: 明确的「答案: xxx」
    ans_patterns = [
        r'\n\s*(?:答案|解答|解析|参考答案)[：:\s]+(.*)',
        r'(?:答案|解答|解析|参考答案)[：:\s]+(.*)',
    ]

    for pat in ans_patterns:
        m = re.search(pat, segment, re.DOTALL)
        if m:
            answer = m.group(1).strip()
            question = segment[:m.start()].strip()
            # 限制答案长度
            answer = answer[:1000]
            return question, answer

    # 模式2: 最后一行是答案（以 "故选" "因此" "答案为" 开头）
    lines = segment.split("\n")
    if len(lines) >= 2:
        last_line = lines[-1].strip()
        if re.match(r'^(故选|因此|所以|答案为|故|综上)', last_line):
            answer = last_line
            question = "\n".join(lines[:-1]).strip()
            return question, answer

    # 模式3: 选择题答案（单独一个字母 A/B/C/D）
    if re.match(r'^[A-Da-d]$', segment):
        return "", segment.upper()

    # 模式4: 选择题答案行（如 "1-5: ABCDD"）
    multi_choice = re.match(r'^[\d\-\s,，]+[：:]\s*([A-Da-d\s]+)$', segment)
    if multi_choice:
        return "", multi_choice.group(1).strip()

    # 模式5: 短文本整体当作答案
    if len(segment) < 80:
        return "", segment

    # 默认：全部当题干，答案为空
    return segment, ""

